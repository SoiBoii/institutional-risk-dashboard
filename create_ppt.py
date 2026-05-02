from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Slide Layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # --- Slide 1: Title Slide ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Advanced Institutional Portfolio Risk & Analytics Platform"
    subtitle.text = "A Cyberpunk-Themed Quantitative Dashboard"

    # Helper function to add slides
    def add_slide(title_text, bullet_points):
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = title_text
        tf = body_shape.text_frame
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.level = 0
            
    # --- Slide 2: Project Overview ---
    add_slide(
        "Project Overview",
        [
            "A production-ready, full-stack web application designed for quantitative portfolio analysis.",
            "Backend Stack: Python, Flask, Pandas, NumPy, SciPy, and yfinance.",
            "Frontend Stack: HTML5, Vanilla JavaScript, Chart.js, and Tailwind CSS.",
            "Purpose: To deliver complex risk calculations, trade generation, and actionable insights in real-time."
        ]
    )

    # --- Slide 3: Frontend Architecture & UI/UX ---
    add_slide(
        "Frontend Architecture & UI/UX",
        [
            "Institutional Cyberpunk Aesthetic: Deep obsidian backgrounds, translucent glass panels, and vibrant neon accents.",
            "Dynamic KPI Cards: Instantly view Expected Return, Sharpe Ratio, Max Drawdown, and Conditional Value at Risk (CVaR).",
            "Universal Deep-Dive Modals: Every chart and metric is clickable, expanding into a full-screen terminal overlay for deeper analysis.",
            "Responsive Grid Layout: Optimizes data visualization across all screen sizes."
        ]
    )

    # --- Slide 4: Quantitative Engine (Backend) ---
    add_slide(
        "Quantitative Engine (Backend)",
        [
            "Robust Data Pipeline: Asynchronously fetches and cleans historical market data via the Yahoo Finance API.",
            "Portfolio Mathematics: Computes daily percentage returns, annualized volatility, and asset correlation matrices.",
            "Optimization Algorithms: Uses SciPy to calculate the Maximum Sharpe Ratio and Minimum Volatility portfolio weights.",
            "Vectorized Operations: Leverages NumPy arrays to ensure sub-second dashboard reactivity even with heavy math."
        ]
    )

    # --- Slide 5: Advanced Analytics ---
    add_slide(
        "Advanced Analytics: Monte Carlo & Factor Exposure",
        [
            "Monte Carlo Efficient Frontier: Generates 1,000+ random portfolio weightings to visualize optimal risk/return boundaries.",
            "Wealth Projection: Utilizes Geometric Brownian Motion (GBM) to forecast 10th, 50th, and 90th percentile wealth paths over 10 years.",
            "Factor Exposure Radar: Runs Multiple Linear Regression to compute portfolio betas against Market, Size, and Value proxies.",
            "Visualized seamlessly on the frontend using dynamic Chart.js radar and scatter plots."
        ]
    )

    # --- Slide 6: Trade Rebalancing & Sentiment Analysis ---
    add_slide(
        "Trade Rebalancing & Sentiment Analysis",
        [
            "Trade Rebalancing Engine: Accepts a 'Total Capital' input and computes the exact dollar flow and shares required to reach optimal weights.",
            "Execution Terminal: Displays the generated BUY (neon green) and SELL (neon red) orders in a monospaced command prompt view.",
            "Alternative Data Integration: Fetches real-time financial news headlines for the targeted assets.",
            "Live Comm-Link: Performs NLP sentiment analysis using TextBlob, scoring headlines as BULL, BEAR, or NEUTRAL."
        ]
    )

    # --- Slide 7: Related Research Mention ---
    add_slide(
        "Related Research: Password Security",
        [
            "Paper Title: Security Performance Tradeoff Modeling of Modern Password Hashing Algorithms.",
            "The paper introduces the Security Efficiency Index (SEI) to evaluate the balance between computational efficiency and attack resistance.",
            "It rigorously analyzes algorithms like SHA-256, bcrypt, PBKDF2, and Argon2.",
            "Conclusion: Argon2 is identified as providing the maximum SEI, making it the most effective algorithm for resisting GPU-accelerated brute-force attacks due to its high memory hardness."
        ]
    )

    # Save the presentation
    output_path = "/Users/saithejas/Desktop/FINAL_PROJECT_MCA/Final_Project_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
